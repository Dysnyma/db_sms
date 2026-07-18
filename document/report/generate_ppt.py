"""答辩 PPT —— 重点：解决问题的思路·方法·手段·成效"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

C1 = RGBColor(0x1A, 0x52, 0x76)  # 深蓝
C2 = RGBColor(0x2E, 0x86, 0xC1)  # 中蓝
C3 = RGBColor(0x4C, 0xAF, 0x50)  # 绿
C4 = RGBColor(0xEB, 0xF5, 0xFB)  # 浅蓝
C5 = RGBColor(0xFF, 0xFF, 0xFF)
C6 = RGBColor(0x2C, 0x3E, 0x50)
C7 = RGBColor(0x7F, 0x8C, 0x8D)
C8 = RGBColor(0xE7, 0x4C, 0x3C)  # 红

def bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def box(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background(); return s

def rbox(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background(); return s

def txt(slide, l, t, w, h, text, sz=14, bold=False, color=C6, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tx.text_frame.word_wrap = True
    p = tx.text_frame.paragraphs[0]; p.text = text; p.alignment = align
    p.runs[0].font.size = Pt(sz); p.runs[0].font.bold = bold
    p.runs[0].font.color.rgb = color; p.runs[0].font.name = 'Microsoft YaHei'
    return tx

def multiline(slide, l, t, w, h, lines, sz=13, color=C6):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tx.text_frame.word_wrap = True
    for i, (text, bd, clr) in enumerate(lines):
        if i == 0: p = tx.text_frame.paragraphs[0]
        else: p = tx.text_frame.add_paragraph()
        p.text = text; p.space_after = Pt(4)
        p.runs[0].font.size = Pt(sz); p.runs[0].font.bold = bd
        p.runs[0].font.color.rgb = clr or color
    return tx

def slide_problem(slide, title, problem, solution_items):
    """标准问题→解决方案布局"""
    box(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), C1)
    txt(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5), title, sz=26, bold=True, color=C1)

    # 问题框
    bar = rbox(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.0), RGBColor(0xFD, 0xED, 0xED))
    txt(slide, Inches(0.7), Inches(1.15), Inches(0.8), Inches(0.4), '❓ 问题', sz=14, bold=True, color=C8)
    txt(slide, Inches(1.5), Inches(1.2), Inches(11), Inches(0.8), problem, sz=13, color=C6)

    # 方案
    items_data = []
    for label, detail in solution_items:
        items_data.append((label, True, C2))
        items_data.append((f'     {detail}', False, C6))

    multiline(slide, Inches(0.7), Inches(2.4), Inches(11.5), Inches(4.5), items_data, sz=13)


# ════════════════════════════════════════
# Slide 1: 封面
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide, C1)
box(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), C3)
txt(slide, Inches(1), Inches(2.0), Inches(11.33), Inches(1.0),
    '数据库系统课程设计', sz=40, bold=True, color=C5, align=PP_ALIGN.CENTER)
txt(slide, Inches(1), Inches(3.0), Inches(11.33), Inches(0.8),
    '学生成绩管理系统', sz=32, color=RGBColor(0x85, 0xC1, 0xE9), align=PP_ALIGN.CENTER)
box(slide, Inches(5.5), Inches(3.9), Inches(2.33), Inches(0.04), C3)
for i, (l, v) in enumerate([('专    业', '计算机科学与技术'), ('学生姓名', '蔡坤灿'), ('指导教师', ''), ('日    期', '2026年7月')]):
    txt(slide, Inches(4.5), Inches(4.2+i*0.5), Inches(4.33), Inches(0.45),
        f'{l}：{v}', sz=16, color=C5, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════
# Slide 2: 目录
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); box(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), C1)
txt(slide, Inches(1), Inches(0.3), Inches(4), Inches(0.6), '目  录', sz=32, bold=True, color=C1)
items = [('01', '系统功能总览'), ('02', '技术方案'), ('03-09', '问题与解决方案（核心）'), ('10', '项目特色'), ('11', '设计总结与不足')]
for i, (n, t) in enumerate(items):
    y = Inches(1.5+i*1.0)
    box(slide, Inches(1), y, Inches(0.7), Inches(0.6), C2)
    txt(slide, Inches(1.05), y+Inches(0.08), Inches(0.6), Inches(0.45), n, sz=18, bold=True, color=C5, align=PP_ALIGN.CENTER)
    txt(slide, Inches(2.0), y+Inches(0.1), Inches(8), Inches(0.45), t, sz=18, color=C6)

# ════════════════════════════════════════
# Slide 3: 系统功能总览
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
box(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), C1)
txt(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5), '01  系统功能总览', sz=26, bold=True, color=C1)

roles = [
    ('🎓  学  生', C2, ['可选课程查询（过滤满员/过期/同课异师）', '选课/退选（触发器实时校验+行级锁防超卖）', '我的成绩（Plotly 柱状图+学籍分/GPA）', '学期均分查询（selectbox 下拉）']),
    ('👨‍🏫  教  师', C3, ['逐条录入成绩（Pydantic 校验）', 'CSV 批量导入（按行反馈中文结果）', '查看课程学生（饼图+排行+未录入提醒）']),
    ('📋  教务管理员', C1, ['数据概览（Plotly 6 种图表）', '6 大实体 CRUD + 专业管理', '排课管理（datetime_input 原生组件）', '选课管理（按学号+学期精准查询）', '成绩统计/明细/教师排行/备份恢复']),
]
for i, (title, color, items) in enumerate(roles):
    x = Inches(0.5+i*4.2)
    rbox(slide, x, Inches(1.2), Inches(3.9), Inches(4.5), C4)
    box(slide, x, Inches(1.2), Inches(3.9), Inches(0.7), color)
    txt(slide, x+Inches(0.2), Inches(1.3), Inches(3.5), Inches(0.5), title, sz=18, bold=True, color=C5, align=PP_ALIGN.CENTER)
    multiline(slide, x+Inches(0.2), Inches(2.1), Inches(3.5), Inches(3.5),
              [(f'✓ {item}', False, C6) for item in items], sz=12)

txt(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(0.4),
    '数据规模：20,000 学生 | 300 教师 | 329 班级 | 1,004 排课 | 2,002,688 选课 | 160+ 专业', sz=12, color=C7)

# ════════════════════════════════════════
# Slide 4: 技术方案
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
box(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), C1)
txt(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5), '02  技术方案', sz=26, bold=True, color=C1)

techs = [
    ('🎨 表现层', C4, C1, ['Streamlit Web（22 页面）', 'Plotly 交互式图表', 'Pydantic 前端校验+max_chars']),
    ('🧠 业务逻辑层', C4, C2, ['9 个存储过程封装业务', '5 个触发器维护数据一致性', '3 个视图简化查询']),
    ('🗄️ 数据层', C4, C3, ['MySQL 8.0 / InnoDB', 'SQLAlchemy QueuePool 连接池', '全部逻辑删除(is_deleted)']),
]
for i, (title, bg_c, c, items) in enumerate(techs):
    x = Inches(0.5+i*4.2)
    rbox(slide, x, Inches(1.2), Inches(3.9), Inches(2.8), bg_c)
    txt(slide, x+Inches(0.2), Inches(1.3), Inches(3.5), Inches(0.5), title, sz=16, bold=True, color=c)
    multiline(slide, x+Inches(0.2), Inches(1.9), Inches(3.5), Inches(1.8),
              [(f'• {item}', False, C6) for item in items], sz=13)

# 六个核心库
libs = [('Python 3.12', '编程语言'), ('MySQL 8.0', '数据库'), ('Streamlit', 'Web 框架'),
        ('PyMySQL+Pool', '数据库驱动'), ('Plotly', '交互图表'), ('Pydantic', '数据校验')]
for i, (n, d) in enumerate(libs):
    x = Inches(0.5+i*2.1)
    rbox(slide, x, Inches(4.3), Inches(1.9), Inches(1.5), C5)
    txt(slide, x+Inches(0.1), Inches(4.5), Inches(1.7), Inches(0.4), n, sz=14, bold=True, color=C1, align=PP_ALIGN.CENTER)
    txt(slide, x+Inches(0.1), Inches(5.0), Inches(1.7), Inches(0.4), d, sz=11, color=C7, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════
# Slides 5-9: 问题与解决方案
# ════════════════════════════════════════

# Slide 5: 排序规则冲突
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_problem(slide, '03  问题：MySQL 排序规则冲突（报错 1267）',
    'MySQL 8.0 默认排序规则 utf8mb4_0900_ai_ci 与数据库的 utf8mb4_unicode_ci 不兼容，\n存储过程中字符串比较时报错 1267「Illegal mix of collations」，所有涉及字符串查询的功能全部崩溃。',
    [('💡 思路', '找到冲突根因，在连接层和查询层双重统一排序规则'),
     ('🔧 方法', 'pymysql 连接参数指定 collation + 存储过程显式加 COLLATE 子句'),
     ('🛠️ 手段', '①DB_URL 加 collation=utf8mb4_unicode_ci ②SP/fn 中字符串比较后加 COLLATE utf8mb4_unicode_ci'),
     ('✅ 成效', '1267 错误完全清零，所有存储过程正常运行')])

# Slide 6: 并发选课超卖 → FOR UPDATE
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_problem(slide, '04  问题：最后一个名额同时被多人选走（并发超卖）',
    '两个学生同时选最后一门课（current=29, max=30），两个事务同时读到 current=29，\n都认为有名额而 INSERT 成功，最终变成 31/30 的数据不一致。',
    [('💡 思路', '让第二个事务等待，读取到最新的 current 值再做判断'),
     ('🔧 方法', '触发器中使用 SELECT ... FOR UPDATE 加排他行锁'),
     ('🛠️ 手段', 'trg_enrollment_before_insert 中加 FOR UPDATE，学生 B 等待直到学生 A COMMIT'),
     ('✅ 成效', 'A 选上最后一个名额，B 收到「选课已满！当前 30/30」异常并 ROLLBACK')])

# Slide 7: 200万数据导入
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_problem(slide, '05  问题：200 万条选课数据导入太慢',
    '200 万条选课 + 5 个触发器，逐条 INSERT + 触发逻辑，预计需要数小时。\n触发器每次都会执行 FOR UPDATE 锁、人数更新和 GPA 重算，200 万×5 次 = 1000 万次触发。',
    [('💡 思路', '先删触发器 → 批量导入 → 手动一次重算，避免逐行触发'),
     ('🔧 方法', 'LOAD DATA LOCAL INFILE（比 INSERT 快 20 倍）+ Python 批量 GPA UPDATE'),
     ('🛠️ 手段', '①Python 预删 5 个触发器 ②LOAD DATA INFILE 导入 CSV ③UPDATE 一次性重算全部学生 GPA ④Python 重建触发器'),
     ('✅ 成效', '导入时间从数小时降至约 30 秒，200 万条 + 20000 人 GPA 重算全部完成')])

# Slide 8: 页面卡顿 → 按学号+学期查询
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_problem(slide, '06  问题：200 万条数据全量加载导致页面卡顿',
    '选课管理页面加载全部 200 万条选课记录，Python 内存消耗巨大、Streamlit 渲染超时、\n下拉退选框 200 万选项浏览器直接卡死。',
    [('💡 思路', '用户只需要查特定学生的选课，不要全部加载'),
     ('🔧 方法', '改成按学期+学号精准查询，退改用输入 ID'),
     ('🛠️ 手段', '①学期 selectbox（轻量 DISTINCT 查询）②输入学号查该生该学期选课 ③退选输入选课 ID'),
     ('✅ 成效', '每次查询只返回几行结果，页面秒开，零延迟')])

# Slide 9: 输入校验 → Pydantic
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_problem(slide, '07  问题：表单无校验导致脏数据入库',
    '用户可输入任意非法数据：成绩填 abc、学号填 999999999999、排课上限填超出 INT 范围的超大数字。\n部分数据到 SQL 层才报错，用户体验差且数据库有风险。',
    [('💡 思路', '在 Streamlit 前端和后端存储过程之间增加一层统一的输入校验'),
     ('🔧 方法', 'Pydantic 模型 + validate_or_error 统一函数 + 前端 max_chars 截断'),
     ('🛠️ 手段', '①创建 StudentCreate/TeacherCreate/GradeRecord 等校验模型 ②学号/工号/成绩格式校验 ③text_input 配置 max_chars ④st.toast 中文错误提示'),
     ('✅ 成效', '脏数据零入库，非法输入在前端被拦截并显示中文报错')])

# Slide 10: JS精度 → text_input + max_chars
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_problem(slide, '08  问题：JavaScript 精度丢失导致校验被绕过',
    '排课上限用 st.number_input，底层依赖 JS Number（64 位浮点）。\n输入 46 位超大数字时超出 JS 安全整数范围（2^53），JS 静默截断→Python 收到正常值→校验通过→SQL 爆掉。',
    [('💡 思路', '绕过 JS Number，用 text_input（字符串）传值，Python 端手动 int() 转换'),
     ('🔧 方法', 'text_input + max_chars=5（最多 5 位数字 0~99999）+ 后端 int() + >99999 检查'),
     ('🛠️ 手段', '①排课上限改为 text_input(max_chars=5) ②按钮中 int(max_s) ③max_s > 99999 拦截'),
     ('✅ 成效', '不管用户输入多少位数字，最多输入 5 位，JS 精度问题彻底解决')])

# ════════════════════════════════════════
# Slide 11: 项目特色
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
box(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), C1)
txt(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5), '09  项目特色', sz=26, bold=True, color=C1)

features = [
    ('🎯 双界面交付', 'CLI + Streamlit Web 双模式运行，\n同一套业务函数复用，代码零冗余'),
    ('🔒 并发安全', 'FOR UPDATE 行级锁 + 事务 + \nUNIQUE 索引三层防护防超卖'),
    ('🚀 大数据量', '200 万选课 30 秒导入，\nLOAD DATA + 删触发器 + 批量 GPA'),
    ('✅ 输入校验', 'Pydantic 全表单校验 +\nst.toast 中文提示 + max_chars 截断'),
    ('📊 交互图表', 'Plotly 饼图/柱状图 10+ 处，\n成绩分布/教师排行一目了然'),
    ('🤖 AI 辅助', 'Claude Code 全流程辅助，\n人工逐行审核+80+次 Git 提交'),
]
for i, (t, d) in enumerate(features):
    x = Inches(0.5+(i%3)*4.2)
    y = Inches(1.1+(i//3)*3.0)
    rbox(slide, x, y, Inches(3.9), Inches(2.7), C4)
    txt(slide, x+Inches(0.2), y+Inches(0.15), Inches(3.5), Inches(0.4), t, sz=16, bold=True, color=C1)
    txt(slide, x+Inches(0.2), y+Inches(0.7), Inches(3.5), Inches(1.8), d, sz=12, color=C6)

# ════════════════════════════════════════
# Slide 12: 设计总结与不足
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
box(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), C1)
txt(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5), '10  设计总结与不足', sz=26, bold=True, color=C1)

# 总结
rbox(slide, Inches(0.5), Inches(1.1), Inches(6.0), Inches(4.0), C4)
txt(slide, Inches(0.7), Inches(1.2), Inches(5.5), Inches(0.4), '✅ 设计总结', sz=18, bold=True, color=C3)
summary = [
    '数据库：8 表 · 3 视图 · 5 触发器 · 9 存储过程 · 2 函数',
    '应用层：CLI + Streamlit 双界面 · 22 个 Web 页面',
    '校验层：Pydantic 全表单校验 · Plotly 交互图表',
    '并发：FOR UPDATE 行锁 · 事务机制 · 唯一索引三层防护',
    '数据：20,000 学生 · 300 教师 · 200 万选课 · 30 秒导入',
    '版本：Git 80+ 次提交 · GitHub 托管 · AI 全程辅助',
]
multiline(slide, Inches(0.7), Inches(1.7), Inches(5.5), Inches(3.2),
          [(f'• {item}', False, C6) for item in summary], sz=12)

# 不足
rbox(slide, Inches(6.8), Inches(1.1), Inches(6.0), Inches(4.0), RGBColor(0xFD, 0xED, 0xED))
txt(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.4), '⚠️ 存在不足', sz=18, bold=True, color=C8)
limits = [
    '缺少自动化单元测试和 CI/CD 流水线',
    '未实现 Redis 缓存和读写分离',
    '物理设计可进一步系统化 EXPLAIN 分析',
    'AI 对 Streamlit 组件行为理解不深入，st.form 迭代 3 轮',
    '大数据量下部分页面查询仍需优化',
]
multiline(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(3.2),
          [(f'• {item}', False, C6) for item in limits], sz=12)

# 致谢
box(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.06), C2)
txt(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.6),
    '谢谢老师！请各位老师批评指正 🙏   课程设计只是起点，工程化思维才是真正的收获', sz=16, bold=True, color=C1, align=PP_ALIGN.CENTER)

# ══ 保存 ══
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), '3.4答辩PPT.pptx')
prs.save(out)
sz = os.path.getsize(out)//1024
print(f'OK: {sz} KB -> {out}')
